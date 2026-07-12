import os
import sys
import subprocess

# 1. Pastikan paket python-docx dan python-pptx terinstall
def ensure_packages():
    packages = {"docx": "python-docx", "pptx": "python-pptx"}
    for mod_name, pip_name in packages.items():
        try:
            __import__(mod_name)
        except ImportError:
            print(f"[INFO] Menginstall paket {pip_name} untuk pembuatan dokumen...")
            subprocess.check_call([sys.executable, "-m", "pip", "install", pip_name, "--quiet"])

ensure_packages()

import docx
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml import OxmlElement
from docx.oxml.ns import qn

import pptx
from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt
from pptx.dml.color import RGBColor as PRGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ==============================================================================
# FUNGSI BANTUAN UNTUK WORD (DOCX)
# ==============================================================================
def set_cell_background(cell, fill_color):
    tcPr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement('w:shd')
    shd.set(qn('w:val'), 'clear')
    shd.set(qn('w:color'), 'auto')
    shd.set(qn('w:fill'), fill_color)
    tcPr.append(shd)

def create_word_document(filename):
    doc = docx.Document()
    
    # Pengaturan Margin
    sections = doc.sections
    for section in sections:
        section.top_margin = Inches(1.0)
        section.bottom_margin = Inches(1.0)
        section.left_margin = Inches(1.0)
        section.right_margin = Inches(1.0)

    # Styling Dasar
    styles = doc.styles
    normal_style = styles['Normal']
    normal_font = normal_style.font
    normal_font.name = 'Calibri'
    normal_font.size = Pt(11)
    normal_font.color.rgb = RGBColor(51, 51, 51) # Dark gray

    # --- HALAMAN JUDUL (COVER) ---
    title_p = doc.add_paragraph()
    title_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title_run = title_p.add_run("OPTIMASI KONTROL CURSOR DAN INPUT KEYBOARD VIRTUAL BERBASIS HAND-LANDMARK DETECTION MENGGUNAKAN ARSITEKTUR MEDIAPIPE SECARA REAL TIME\n(DACTYL-AI)")
    title_run.font.name = 'Arial'
    title_run.font.size = Pt(18)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(0, 51, 102) # Deep Blue
    
    sub_p = doc.add_paragraph()
    sub_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sub_run = sub_p.add_run("Dokumentasi Teknis dan Analisis Komprehensif Sistem V2.5\nMagic Cursor Bimanual System")
    sub_run.font.size = Pt(13)
    sub_run.font.italic = True
    sub_run.font.color.rgb = RGBColor(102, 102, 102)

    doc.add_paragraph().add_run("\n" * 2)

    # Box Spesifikasi Pengetahuan
    table_spec = doc.add_table(rows=3, cols=2)
    table_spec.alignment = WD_TABLE_ALIGNMENT.CENTER
    table_spec.autofit = False
    
    specs = [
        ("Judul Sistem", "Optimasi Kontrol Cursor dan Input Keyboard Virtual Berbasis Hand-Landmark Detection Menggunakan Arsitektur MediaPipe Secara Real Time (Dactyl-AI) V2.5"),
        ("Sumber Pengetahuan", "MediaPipe Hand Tracking, OpenCV (Real-Time Camera Input), Scikit-Learn (SVM & Scaler), PyAutoGUI & Win32 API, Rule-Based Gesture Logic, Geometry & Spatial Mapping."),
        ("Pengetahuan AI & ML", "• 126 atribut koordinat spasial 3D (X, Y, Z) dari 21 titik hand-landmarks (Tangan Kiri & Kanan)\n• Jarak Euclidean (Euclidean distance) antar jari & Sudut rotasi pergelangan tangan (atan2)\n• Kecepatan geser (velocity dx/dt), Bounding Box area, & Kalkulasi temporal (durasi penahanan gestur)\n• Support Vector Machine (SVM) dengan Kernel RBF & Normalisasi Data (StandardScaler)\n• Exponential Moving Average (EMA) untuk Heuristic Jitter Filtering\n• Rule-Based Geometrical Classification")
    ]
    
    col_widths = [Inches(1.8), Inches(4.7)]
    for idx, (label, val) in enumerate(specs):
        row = table_spec.rows[idx]
        cell_lbl = row.cells[0]
        cell_val = row.cells[1]
        
        cell_lbl.width = col_widths[0]
        cell_val.width = col_widths[1]
        
        p0 = cell_lbl.paragraphs[0]
        r0 = p0.add_run(label)
        r0.bold = True
        r0.font.color.rgb = RGBColor(0, 51, 102)
        set_cell_background(cell_lbl, "E6F0FA")
        
        p1 = cell_val.paragraphs[0]
        p1.add_run(val)
        set_cell_background(cell_val, "F9FBFD")
        
    doc.add_page_break()

    # --- BAB 1: PENDAHULUAN & RINGKASAN EKSEKUTIF ---
    h1 = doc.add_heading("1. RINGKASAN EKSEKUTIF (EXECUTIVE SUMMARY)", level=1)
    h1.style.font.color.rgb = RGBColor(0, 51, 102)
    
    p = doc.add_paragraph(
        "Dactyl-AI adalah inovasi sistem Human-Computer Interaction (HCI) tanpa sentuhan (touchless interface) "
        "yang memungkinkan pengguna mengontrol kursor mouse, mengetik pada keyboard virtual transparan, mengatur multimedia "
        "dan kecerahan layar, hingga melakukan navigasi multitasking sistem operasi hanya menggunakan gestur tangan di depan "
        "kamera web (webcam) secara real-time."
    )
    p = doc.add_paragraph(
        "Sistem ini memecahkan dua masalah utama dalam interaksi berbasis visi komputer konvensional:\n"
        "1. Ketergantungan pada Perangkat Keras Fisik: Menggantikan mouse dan keyboard fisik dengan antarmuka optik yang ergonomis.\n"
        "2. Oklusi dan Konflik Gestur: Mengatasi masalah jari saling menutupi (occlusion) dengan menerapkan konsep 'Dual-Wielding / Bimanual Interface', "
        "di mana tangan kanan dikhususkan untuk navigasi dan klik, sedangkan tangan kiri dikhususkan untuk shortcut, kontrol sistem, dan penunjang keyboard."
    )

    # --- BAB 2: ARSITEKTUR & TEKNOLOGI PENDUKUNG ---
    h1 = doc.add_heading("2. ARSITEKTUR SISTEM & TEKNOLOGI PENDUKUNG", level=1)
    h1.style.font.color.rgb = RGBColor(0, 51, 102)
    
    p = doc.add_paragraph(
        "Sistem Dactyl-AI dibangun di atas ekosistem Python 3.8+ dengan arsitektur modular yang mengintegrasikan "
        "enam teknologi inti sebagai sumber pengetahuannya:"
    )
    
    tech_list = [
        ("MediaPipe Hands (Google)", "Berfungsi sebagai feature extractor spasial yang mendeteksi hingga 2 tangan secara simultan dan mengekstrak 21 titik sendi (hand-landmarks) koordinat spasial 3D (X, Y, Z) per tangan."),
        ("Scikit-Learn (SVM & Scaler)", "Mesin kecerdasan buatan (AI Brain) yang mengklasifikasikan vektor 126 dimensi ke dalam 8 kelas gestur menggunakan Support Vector Machine (SVM) dengan Kernel RBF."),
        ("OpenCV (Computer Vision)", "Menangani akuisisi video secara real-time dari webcam (1280x720), mirroring frame, rendering antarmuka visual (hologram UI, progress bar, bingkai kamera), serta manipulasi citra."),
        ("PyAutoGUI", "Mengonversi instruksi gestur menjadi event sistem operasi seperti gerakan kursor, klik kiri/kanan, hold-drag, scroll vertikal/horizontal, serta keystroke shortcut keyboard."),
        ("Win32 API (ctypes & wintypes)", "Mengatur window aplikasi agar bersifat Always on Top (Topmost), semi-transparan (60% opacity), dan memosisikan diri di pojok kanan bawah monitor agar tidak mengganggu area kerja."),
        ("PowerShell Asynchronous", "Digunakan untuk pengontrolan kecerahan monitor (WmiSetBrightness) di latar belakang (background thread) tanpa menyebabkan lag atau bottleneck pada video feed.")
    ]
    for tech, desc in tech_list:
        bp = doc.add_paragraph(style='List Bullet')
        r_tech = bp.add_run(f"{tech}: ")
        r_tech.bold = True
        r_tech.font.color.rgb = RGBColor(0, 51, 102)
        bp.add_run(desc)

    # --- BAB 3: PENGETAHUAN AI & PIPELINE REAL-TIME ---
    h1 = doc.add_heading("3. PENGETAHUAN AI & ALUR KERJA SISTEM (REAL-TIME PIPELINE)", level=1)
    h1.style.font.color.rgb = RGBColor(0, 51, 102)
    
    p = doc.add_paragraph(
        "Keunggulan utama Dactyl-AI terletak pada penggabungan pemrosesan Machine Learning numerik dan logika geometri spasial. "
        "Setiap frame video diproses melalui pipeline 5 tahap dengan frekuensi 24-30 FPS:"
    )
    
    pipeline_steps = [
        ("Tahap 1: Akuisisi & Preprocessing Video", "Webcam menangkap frame video 1280x720, melakukan pencerminan (cv2.flip) agar interaksi terasa alami seperti berkaca, lalu mengonversi ruang warna dari BGR ke RGB untuk konsumsi MediaPipe."),
        ("Tahap 2: Ekstraksi Fitur Spasial 3D (126 Atribut)", "MediaPipe mendeteksi 21 koordinat (X, Y, Z) per tangan. Sistem menyusun vektor 126 dimensi (63 atribut Tangan Kiri + 63 atribut Tangan Kanan). Jika satu tangan tidak terdeteksi, atributnya diisi nilai 0.0. Sistem juga melakukan koreksi label cermin otomatis."),
        ("Tahap 3: Normalisasi & Klasifikasi SVM RBF", "Vektor 126 dimensi dinormalisasi menggunakan StandardScaler agar rentang nilainya seragam dengan data latih. Model Support Vector Machine (SVM) dengan Kernel RBF memprediksi probabilitas label gestur (Kelas 0 hingga 7)."),
        ("Tahap 4: Eksekusi Logika Geometri & Temporal", "Berdasarkan label prediksi, sistem menghitung Jarak Euclidean antar jari, sudut rotasi pergelangan tangan (atan2), kecepatan geser (dx/dt), dan durasi penahanan gestur (timer temporal)."),
        ("Tahap 5: Rendering Visual & UI Feedback", "Menampilkan overlay hologram, titik hover keyboard, progress bar interaktif, dan status gestur di atas layar kamera transparan.")
    ]
    for step, desc in pipeline_steps:
        p_step = doc.add_paragraph()
        r_step = p_step.add_run(f"• {step}\n")
        r_step.bold = True
        r_step.font.color.rgb = RGBColor(0, 51, 102)
        p_step.add_run(desc)

    # --- BAB 4: DAFTAR KELAS GESTUR ---
    h1 = doc.add_heading("4. DAFTAR KELAS GESTUR & MEKANISME KERJA (LABEL 0 - 7)", level=1)
    h1.style.font.color.rgb = RGBColor(0, 51, 102)
    
    p = doc.add_paragraph("Sistem mengklasifikasikan 8 label gestur dengan aturan logika geometri spasial sebagai berikut:")
    
    table_gestures = doc.add_table(rows=1, cols=4)
    table_gestures.style = 'Table Grid'
    table_gestures.alignment = WD_TABLE_ALIGNMENT.CENTER
    
    headers = ["Label", "Nama Gestur", "Fungsi Kontrol", "Mekanisme Kerja & Geometri AI"]
    hdr_cells = table_gestures.rows[0].cells
    for i, header_text in enumerate(headers):
        p_hdr = hdr_cells[i].paragraphs[0]
        r_hdr = p_hdr.add_run(header_text)
        r_hdr.bold = True
        r_hdr.font.color.rgb = RGBColor(255, 255, 255)
        set_cell_background(hdr_cells[i], "003366")
        
    gestures_data = [
        ("Label 0", "Telunjuk Lurus\n(Navigation)", "Menggerakkan kursor mouse di layar", "Menggunakan koordinat ujung telunjuk kanan (Landmark 8). Dilengkapi Exponential Moving Average (EMA) alpha=0.25 untuk Heuristic Jitter Filtering penghilang tremor."),
        ("Label 1", "Cubit / Pinch\n(Click & Drag)", "Klik Kiri, Klik Kanan, Hold-Drag", "Jarak Euclidean jempol (L4) ke telunjuk (L8) < 0.045. Klik kiri (<0.3s), Hold-Drag (>0.3s), Klik kanan (Cubit jempol + jari tengah L12 dengan cooldown 0.4s)."),
        ("Label 2", "C-Shape / Bola\n(Multimedia)", "Volume Audio & Brightness Layar", "Sudut rotasi pergelangan (L0) ke jari tengah (L9) menggunakan atan2. Tangan Kanan untuk Volume, Tangan Kiri untuk Brightness via PowerShell Asynchronous."),
        ("Label 3", "Peace / V-Sign\n(Keyboard)", "Memunculkan Keyboard Virtual", "Tahan gestur 2 detik untuk memunculkan keyboard 3 baris QWERTY. Mengetik dengan hovering telunjuk selama 1.0 detik (dilengkapi progress bar hijau)."),
        ("Label 4", "Frame Kamera\n(Screenshot)", "Tangkapan Layar Custom Area", "Membentuk bingkai kotak dari jempol & telunjuk kedua tangan. Kalkulasi Bounding Box area ditahan 1.5 detik, disertai efek visual flash kamera."),
        ("Label 5", "Dua Jari Ninja\n(Scroll Mode)", "Scroll Vertikal & Horizontal", "Pelacakan delta posisi (dx, dy) jari tengah kanan (L9). Pergerakan vertikal memicu scroll, horizontal memicu hscroll dengan akselerasi responsif."),
        ("Label 6", "Tangan Terbuka\n(Task Switcher)", "Alt+Tab & Alt+Shift+Tab", "Eksklusif Tangan Kiri. Menghitung kecepatan geser (velocity = dx/dt). Swipe cepat kanan (>1.3) Next Tab, kiri (<-1.3) Prev Tab, cooldown 1.2s."),
        ("Label 7", "Fist / Mengepal\n(Standby Lock)", "Mengunci / Membuka Sistem", "Genggam tangan selama 1.0 detik untuk Lock sistem (semua kontrol dimatikan + overlay merah). Genggam lagi 1.0 detik untuk Unlock.")
    ]
    
    col_w = [Inches(0.9), Inches(1.5), Inches(1.8), Inches(2.3)]
    for row_idx, data in enumerate(gestures_data):
        row_cells = table_gestures.add_row().cells
        for col_idx, text in enumerate(data):
            row_cells[col_idx].width = col_w[col_idx]
            p_cell = row_cells[col_idx].paragraphs[0]
            p_cell.add_run(text)
            if row_idx % 2 == 1:
                set_cell_background(row_cells[col_idx], "F2F6FA")

    doc.add_page_break()

    # --- BAB 5: KEUNGGULAN TEKNIS & STRUKTUR KODE ---
    h1 = doc.add_heading("5. KEUNGGULAN TEKNIS & STRUKTUR KODE PROGRAM", level=1)
    h1.style.font.color.rgb = RGBColor(0, 51, 102)
    
    p = doc.add_paragraph("Dactyl-AI memiliki 5 keunggulan arsitektural utama yang membedakannya dari sistem serupa:")
    
    adv_list = [
        ("Dual-Wielding Bimanual Task Division", "Pembagian tugas eksklusif mencegah oklusi kamera dan memungkinkan interaksi kompleks secara simultan."),
        ("Heuristic Jitter Filtering (EMA Alpha=0.25)", "Penyaringan koordinat numerik meredam getaran alami tangan sehingga pergerakan kursor sepresisi mouse optik."),
        ("Lightweight AI Engine (1D Vector vs 2D Image)", "Klasifikasi SVM pada vektor numerik 126 dimensi membutuhkan komputasi sangat ringan dibanding model CNN/Deep Learning citra gambar."),
        ("Non-Intrusive Always-on-Top GUI", "Manipulasi Win32 API membuat jendela kamera selalu berada di atas (topmost), transparan 60%, dan berada di pojok layar."),
        ("Asynchronous PowerShell Threading", "Eksekusi perintah WMI Windows di background thread menjaga FPS kamera tetap stabil di 24-30 FPS.")
    ]
    for adv, desc in adv_list:
        bp = doc.add_paragraph(style='List Bullet')
        r_adv = bp.add_run(f"{adv}: ")
        r_adv.bold = True
        r_adv.font.color.rgb = RGBColor(0, 51, 102)
        bp.add_run(desc)

    h2 = doc.add_heading("Struktur Modul Program (Folder src/)", level=2)
    h2.style.font.color.rgb = RGBColor(0, 51, 102)
    
    p = doc.add_paragraph(
        "• data_collection.py : Perekaman dataset bimanual. Mengambil 126 koordinat dari MediaPipe dan menyimpannya ke CSV dengan kontrol tombol [0]-[7].\n"
        "• train_model.py : Membagi data 80% Train, 10% Validation, 10% Test dengan stratified sampling. Melakukan normalisasi StandardScaler dan melatih Support Vector Machine (SVM RBF, C=10). Menyimpannya ke models/svm_gesture_model.pkl.\n"
        "• run_system.py : Modul eksekusi utama yang memuat model AI dan scaler, menjalankan loop kamera real-time, mendeteksi landmark, serta mengeksekusi logika kontrol kursor dan UI hologram."
    )

    # --- BAB 6: KESIMPULAN ---
    h1 = doc.add_heading("6. KESIMPULAN", level=1)
    h1.style.font.color.rgb = RGBColor(0, 51, 102)
    
    p = doc.add_paragraph(
        "Dactyl-AI berhasil menggabungkan presisi ekstraksi koordinat MediaPipe, kehandalan klasifikasi Support Vector Machine (SVM), "
        "serta kecerdasan logika geometri spasial untuk menciptakan sistem kontrol antarmuka masa depan yang sepenuhnya tanpa sentuhan. "
        "Dengan performa real-time (24-30 FPS), fitur peredam getaran (jitter filter), dan desain antarmuka transparan yang ergonomis, "
        "sistem ini menjadi solusi Human-Computer Interaction yang sangat praktis, stabil, dan futuristik."
    )

    doc.save(filename)
    print(f"[BERHASIL] Dokumen Word disimpan di: {os.path.abspath(filename)}")


# ==============================================================================
# FUNGSI BANTUAN UNTUK POWERPOINT (PPTX)
# ==============================================================================
def create_powerpoint_presentation(filename):
    prs = Presentation()
    prs.slide_width = PInches(13.333)
    prs.slide_height = PInches(7.5)

    blank_layout = prs.slide_layouts[6] # Blank slide layout

    def add_header(slide, title_text, category_text="DACTYL-AI V2.5 - SYSTEM ARCHITECTURE"):
        # Header banner
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, PInches(13.333), PInches(1.1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = PRGBColor(0, 51, 102) # Deep Blue
        shape.line.fill.background()
        
        # Category / Top small text
        txBox = slide.shapes.add_textbox(PInches(0.6), PInches(0.12), PInches(12), PInches(0.3))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = category_text.upper()
        p.font.size = PPt(11)
        p.font.bold = True
        p.font.color.rgb = PRGBColor(176, 196, 222) # Light steel blue
        
        # Title text
        txBox2 = slide.shapes.add_textbox(PInches(0.6), PInches(0.35), PInches(12), PInches(0.6))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = title_text
        p2.font.size = PPt(24)
        p2.font.bold = True
        p2.font.color.rgb = PRGBColor(255, 255, 255)

    def add_footer(slide, slide_num):
        txBox = slide.shapes.add_textbox(PInches(0.6), PInches(7.0), PInches(12), PInches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"Dactyl-AI V2.5: Real-Time Hand-Landmark Control System  |  Slide {slide_num}"
        p.font.size = PPt(10)
        p.font.color.rgb = PRGBColor(128, 128, 128)

    # --------------------------------------------------------------------------
    # SLIDE 1: COVER
    # --------------------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, PInches(13.333), PInches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = PRGBColor(0, 40, 80)
    bg1.line.fill.background()
    
    # Accent bar
    bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, PInches(1.0), PInches(1.8), PInches(0.15), PInches(3.8))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRGBColor(0, 165, 255) # Bright Blue/Cyan
    bar.line.fill.background()

    txBox = slide1.shapes.add_textbox(PInches(1.4), PInches(1.8), PInches(11), PInches(3.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "OPTIMASI KONTROL CURSOR DAN INPUT KEYBOARD VIRTUAL BERBASIS HAND-LANDMARK DETECTION"
    p.font.size = PPt(28)
    p.font.bold = True
    p.font.color.rgb = PRGBColor(255, 255, 255)
    
    p2 = tf.add_paragraph()
    p2.text = "Menggunakan Arsitektur MediaPipe Secara Real Time (Dactyl-AI V2.5)"
    p2.font.size = PPt(22)
    p2.font.color.rgb = PRGBColor(0, 180, 255)
    p2.space_before = PPt(15)
    
    p3 = tf.add_paragraph()
    p3.text = "Magic Cursor Bimanual System  |  Support Vector Machine & Spatial Geometry"
    p3.font.size = PPt(14)
    p3.font.italic = True
    p3.font.color.rgb = PRGBColor(200, 220, 240)
    p3.space_before = PPt(30)

    # --------------------------------------------------------------------------
    # SLIDE 2: RINGKASAN EKSEKUTIF
    # --------------------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    add_header(slide2, "Ringkasan Eksekutif (Executive Summary)")
    add_footer(slide2, 2)
    
    txBox = slide2.shapes.add_textbox(PInches(0.8), PInches(1.5), PInches(11.7), PInches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    items = [
        ("Inovasi Touchless HCI", "Dactyl-AI adalah sistem Human-Computer Interaction tanpa sentuhan untuk mengontrol kursor, mengetik keyboard virtual, multimedia, & navigasi OS hanya dengan gestur tangan di depan webcam."),
        ("Mengatasi Ketergantungan Hardware", "Menggantikan mouse dan keyboard fisik dengan antarmuka optik transparan yang ergonomis dan tangguh."),
        ("Solusi Masalah Oklusi (Jari Saling Menutupi)", "Menerapkan konsep 'Dual-Wielding / Bimanual Interface' yang membagi tugas secara eksklusif antara tangan kanan (navigasi & klik) dan tangan kiri (shortcut & kontrol sistem)."),
        ("Kombinasi AI & Geometri Spasial", "Diperkuat oleh integrasi Computer Vision (MediaPipe), Machine Learning (SVM Kernel RBF), dan Rule-Based Geometrical Logic secara real time pada 24-30 FPS.")
    ]
    for idx, (head, body) in enumerate(items):
        p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
        r1 = p.add_run()
        r1.text = f"• {head}: "
        r1.font.bold = True
        r1.font.size = PPt(18)
        r1.font.color.rgb = PRGBColor(0, 51, 102)
        
        r2 = p.add_run()
        r2.text = body
        r2.font.size = PPt(16)
        r2.font.color.rgb = PRGBColor(60, 60, 60)
        p.space_after = PPt(15)

    # --------------------------------------------------------------------------
    # SLIDE 3: SUMBER PENGETAHUAN & TEKNOLOGI PENDUKUNG
    # --------------------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    add_header(slide3, "Sumber Pengetahuan & Teknologi Pendukung")
    add_footer(slide3, 3)
    
    techs = [
        ("MediaPipe Hands", "Pemindai koordinat spasial 3D (X, Y, Z) dari 21 titik sendi untuk 2 tangan simultan."),
        ("Scikit-Learn (SVM & Scaler)", "Mesin AI Brain pengklasifikasi 126 atribut vektor ke dalam 8 kelas gestur."),
        ("OpenCV", "Akuisisi video real-time 1280x720, mirroring, dan rendering UI hologram transparan."),
        ("PyAutoGUI", "Eksekutor event OS: gerakan mouse, klik, hold-drag, scroll, dan keystroke keyboard."),
        ("Win32 API (ctypes)", "Mengatur window Always-on-Top (Topmost), opacity 60%, di pojok kanan bawah monitor."),
        ("PowerShell Asynchronous", "Pengontrol kecerahan layar di background thread tanpa menyebabkan lag pada video feed.")
    ]
    
    for i, (title, desc) in enumerate(techs):
        col = i % 2
        row = i // 2
        left = PInches(0.8 + col * 5.9)
        top = PInches(1.6 + row * 1.7)
        
        box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, PInches(5.6), PInches(1.5))
        box.fill.solid()
        box.fill.fore_color.rgb = PRGBColor(240, 246, 252)
        box.line.color.rgb = PRGBColor(0, 102, 204)
        box.line.width = PPt(1.5)
        
        tf = box.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.bold = True
        p1.font.size = PPt(16)
        p1.font.color.rgb = PRGBColor(0, 51, 102)
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = PPt(13)
        p2.font.color.rgb = PRGBColor(70, 70, 70)
        p2.space_before = PPt(5)

    # --------------------------------------------------------------------------
    # SLIDE 4: PENGETAHUAN AI & MACHINE LEARNING
    # --------------------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    add_header(slide4, "Pengetahuan AI & Machine Learning Dactyl-AI")
    add_footer(slide4, 4)
    
    txBox = slide4.shapes.add_textbox(PInches(0.8), PInches(1.5), PInches(11.7), PInches(5.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    ai_points = [
        ("126 Atribut Koordinat Spasial 3D", "Ekstraksi 21 titik hand-landmarks (X, Y, Z) dari Tangan Kiri (63) dan Tangan Kanan (63) menjadi vektor numerik 1D terintegrasi."),
        ("Kalkulasi Geometri Spasial", "Memanfaatkan Jarak Euclidean (Euclidean distance) antar jari untuk deteksi cubitan, sudut rotasi pergelangan tangan (atan2) untuk kontrol volume/brightness, dan Bounding Box area untuk screenshot."),
        ("Pelacakan Dinamis & Temporal", "Perhitungan kecepatan geser (velocity dx/dt) untuk swipe tab, serta kalkulasi temporal (timer penahanan gestur 1.0s - 2.0s) untuk trigger keyboard & lock."),
        ("Support Vector Machine (SVM RBF)", "Klasifikasi gestur menggunakan kernel RBF yang dilatih pada data ber-StandardScaler, memberikan akurasi tinggi dengan beban komputasi sangat ringan."),
        ("Heuristic Jitter Filtering", "Menerapkan Exponential Moving Average (EMA) dengan alpha = 0.25 untuk menyaring tremor alami tangan sehingga kursor bergerak ultra-mulus.")
    ]
    for idx, (head, body) in enumerate(ai_points):
        p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
        r1 = p.add_run()
        r1.text = f"• {head}: "
        r1.font.bold = True
        r1.font.size = PPt(17)
        r1.font.color.rgb = PRGBColor(0, 51, 102)
        
        r2 = p.add_run()
        r2.text = body
        r2.font.size = PPt(15)
        r2.font.color.rgb = PRGBColor(60, 60, 60)
        p.space_after = PPt(12)

    # --------------------------------------------------------------------------
    # SLIDE 5: ALUR KERJA SISTEM (REAL-TIME PIPELINE)
    # --------------------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    add_header(slide5, "Alur Kerja Sistem (5-Stage Real-Time Pipeline)")
    add_footer(slide5, 5)
    
    steps = [
        ("1. Akuisisi Video", "Webcam 1280x720\nMirroring (cv2.flip)\nBGR to RGB"),
        ("2. Ekstraksi Fitur", "MediaPipe Hands\n21 titik (X,Y,Z) x 2\n126 Vektor Dimensi"),
        ("3. AI Brain (SVM)", "StandardScaler Norm\nSVM Kernel RBF\nPrediksi Label 0-7"),
        ("4. Logika Kontrol", "Euclidean Distance\nTemporal Timer\nPyAutoGUI & Win32"),
        ("5. Rendering UI", "Hologram Overlay\nHover Progress Bar\n24-30 FPS Feed")
    ]
    for i, (title, desc) in enumerate(steps):
        left = PInches(0.5 + i * 2.5)
        top = PInches(2.5)
        
        box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, PInches(2.3), PInches(3.2))
        box.fill.solid()
        box.fill.fore_color.rgb = PRGBColor(0, 51, 102) if i % 2 == 0 else PRGBColor(0, 102, 204)
        box.line.fill.background()
        
        tf = box.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.bold = True
        p1.font.size = PPt(16)
        p1.font.color.rgb = PRGBColor(255, 255, 255)
        p1.alignment = PP_ALIGN.CENTER
        
        p2 = tf.add_paragraph()
        p2.text = "\n" + desc
        p2.font.size = PPt(14)
        p2.font.color.rgb = PRGBColor(220, 235, 255)
        p2.alignment = PP_ALIGN.CENTER

    # --------------------------------------------------------------------------
    # SLIDE 6: BIMANUAL INTERFACE / DUAL-WIELDING
    # --------------------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    add_header(slide6, "Konsep Bimanual Interface (Dual-Wielding)")
    add_footer(slide6, 6)
    
    # Left Hand Box
    boxL = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PInches(0.8), PInches(1.8), PInches(5.6), PInches(4.6))
    boxL.fill.solid()
    boxL.fill.fore_color.rgb = PRGBColor(245, 248, 252)
    boxL.line.color.rgb = PRGBColor(0, 102, 204)
    boxL.line.width = PPt(2)
    
    tfL = boxL.text_frame
    tfL.word_wrap = True
    p = tfL.paragraphs[0]
    p.text = "TANGAN KIRI (System & Shortcut Controller)"
    p.font.bold = True
    p.font.size = PPt(18)
    p.font.color.rgb = PRGBColor(0, 51, 102)
    
    left_tasks = [
        "Brightness Control: Rotasi C-shape (atan2) dengan eksekusi PowerShell Asynchronous.",
        "Task Switcher (Alt+Tab): Swipe horizontal cepat kanan (>1.3) atau kiri (<-1.3).",
        "Keyboard Summon Support: Menahan gestur Peace (V-sign) bersama tangan kanan.",
        "Screenshot Frame Support: Membentuk sudut kiri bingkai kamera bounding box."
    ]
    for t in left_tasks:
        p = tfL.add_paragraph()
        p.text = "• " + t
        p.font.size = PPt(14)
        p.font.color.rgb = PRGBColor(60, 60, 60)
        p.space_before = PPt(10)
        
    # Right Hand Box
    boxR = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PInches(6.9), PInches(1.8), PInches(5.6), PInches(4.6))
    boxR.fill.solid()
    boxR.fill.fore_color.rgb = PRGBColor(245, 248, 252)
    boxR.line.color.rgb = PRGBColor(0, 153, 76)
    boxR.line.width = PPt(2)
    
    tfR = boxR.text_frame
    tfR.word_wrap = True
    p = tfR.paragraphs[0]
    p.text = "TANGAN KANAN (Navigation & Action Controller)"
    p.font.bold = True
    p.font.size = PPt(18)
    p.font.color.rgb = PRGBColor(0, 153, 76)
    
    right_tasks = [
        "Cursor Navigation (Label 0): Telunjuk lurus mengontrol kursor dengan EMA Jitter Filter.",
        "Click & Drag (Label 1): Cubit jempol & telunjuk (<0.045) untuk Klik Kiri / Hold-Drag.",
        "Right Click: Cubit jempol & jari tengah kanan (cooldown 0.4s).",
        "Volume Audio Control: Rotasi pergelangan tangan C-shape searah/berlawanan jarum jam.",
        "Scroll Ninja (Label 5): Delta posisi dx, dy jari tengah untuk scroll vertikal/horizontal."
    ]
    for t in right_tasks:
        p = tfR.add_paragraph()
        p.text = "• " + t
        p.font.size = PPt(14)
        p.font.color.rgb = PRGBColor(60, 60, 60)
        p.space_before = PPt(10)

    # --------------------------------------------------------------------------
    # SLIDE 7: GESTUR NAVIGASI & KLIK (LABEL 0 & 1)
    # --------------------------------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    add_header(slide7, "Gestur Navigasi & Klik (Label 0 & Label 1)")
    add_footer(slide7, 7)
    
    box1 = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, PInches(0.8), PInches(1.6), PInches(11.7), PInches(2.5))
    box1.fill.solid()
    box1.fill.fore_color.rgb = PRGBColor(240, 246, 252)
    box1.line.color.rgb = PRGBColor(0, 102, 204)
    tf1 = box1.text_frame
    tf1.word_wrap = True
    p = tf1.paragraphs[0]
    p.text = "LABEL 0 : TELUNJUK LURUS (Move / Navigation Mode)"
    p.font.bold = True
    p.font.size = PPt(18)
    p.font.color.rgb = PRGBColor(0, 51, 102)
    p_sub = tf1.add_paragraph()
    p_sub.text = "• Fungsi: Menggerakkan kursor mouse secara real-time berdasarkan koordinat ujung telunjuk kanan (Landmark 8).\n• Optimasi AI: Dilengkapi algoritma Heuristic Jitter Filtering menggunakan Exponential Moving Average (EMA) dengan alpha = 0.25. Filter ini meredam tremor alami tangan sehingga kursor bergerak sangat mulus seperti mouse optik fisik."
    p_sub.font.size = PPt(15)
    p_sub.font.color.rgb = PRGBColor(50, 50, 50)
    p_sub.space_before = PPt(8)
    
    box2 = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, PInches(0.8), PInches(4.4), PInches(11.7), PInches(2.6))
    box2.fill.solid()
    box2.fill.fore_color.rgb = PRGBColor(240, 246, 252)
    box2.line.color.rgb = PRGBColor(0, 102, 204)
    tf2 = box2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "LABEL 1 : CUBIT / PINCH (Click & Drag Mode)"
    p.font.bold = True
    p.font.size = PPt(18)
    p.font.color.rgb = PRGBColor(0, 51, 102)
    p_sub = tf2.add_paragraph()
    p_sub.text = "• Kalkulasi Geometri: Menghitung Jarak Euclidean antara ujung jempol (Landmark 4) dan ujung jari lainnya (Threshold < 0.045).\n• Klik Kiri: Cubit jempol + telunjuk dilepas cepat (< 0.3 detik).\n• Hold-Drag: Cubit jempol + telunjuk ditahan (> 0.3 detik) memicu event mouseDown untuk menyeret jendela/item.\n• Klik Kanan: Cubit jempol + jari tengah kanan (dengan cooldown 0.4 detik)."
    p_sub.font.size = PPt(15)
    p_sub.font.color.rgb = PRGBColor(50, 50, 50)
    p_sub.space_before = PPt(8)

    # --------------------------------------------------------------------------
    # SLIDE 8: GESTUR MULTIMEDIA & BRIGHTNESS (LABEL 2 & 5)
    # --------------------------------------------------------------------------
    slide8 = prs.slides.add_slide(blank_layout)
    add_header(slide8, "Kontrol Multimedia & Scroll Ninja (Label 2 & 5)")
    add_footer(slide8, 8)
    
    box1 = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, PInches(0.8), PInches(1.6), PInches(11.7), PInches(2.5))
    box1.fill.solid()
    box1.fill.fore_color.rgb = PRGBColor(240, 246, 252)
    box1.line.color.rgb = PRGBColor(0, 102, 204)
    tf1 = box1.text_frame
    tf1.word_wrap = True
    p = tf1.paragraphs[0]
    p.text = "LABEL 2 : MEGANG BOLA / C-SHAPE (Multimedia & Brightness Control)"
    p.font.bold = True
    p.font.size = PPt(18)
    p.font.color.rgb = PRGBColor(0, 51, 102)
    p_sub = tf1.add_paragraph()
    p_sub.text = "• Kalkulasi Sudut: Menggunakan fungsi matematika atan2 untuk menghitung rotasi pergelangan (Landmark 0) ke jari tengah (Landmark 9).\n• Volume Control (Tangan Kanan): Memutar tangan searah/berlawanan jarum jam mengubah volume audio sistem.\n• Brightness Control (Tangan Kiri): Memutar tangan kiri memicu thread PowerShell Asynchronous (WmiSetBrightness) untuk mencerahkan/meredupkan layar tanpa menyebabkan bottleneck atau lag pada video feed."
    p_sub.font.size = PPt(15)
    p_sub.font.color.rgb = PRGBColor(50, 50, 50)
    p_sub.space_before = PPt(8)

    box2 = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, PInches(0.8), PInches(4.4), PInches(11.7), PInches(2.5))
    box2.fill.solid()
    box2.fill.fore_color.rgb = PRGBColor(240, 246, 252)
    box2.line.color.rgb = PRGBColor(0, 102, 204)
    tf2 = box2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "LABEL 5 : DUA JARI LURUS NINJA (Scroll Mode)"
    p.font.bold = True
    p.font.size = PPt(18)
    p.font.color.rgb = PRGBColor(0, 51, 102)
    p_sub = tf2.add_paragraph()
    p_sub.text = "• Pelacakan Delta Posisi: Posisi jari tengah (Landmark 9) dilacak pergerakan delta-nya (dx, dy) antar frame.\n• Scroll Vertikal & Horizontal: Pergerakan vertikal memicu pyautogui.scroll (naik/turun), sedangkan pergerakan horizontal memicu pyautogui.hscroll (kiri/kanan) dengan faktor akselerasi yang sangat responsif."
    p_sub.font.size = PPt(15)
    p_sub.font.color.rgb = PRGBColor(50, 50, 50)
    p_sub.space_before = PPt(8)

    # --------------------------------------------------------------------------
    # SLIDE 9: KEYBOARD VIRTUAL & SCREENSHOT (LABEL 3 & 4)
    # --------------------------------------------------------------------------
    slide9 = prs.slides.add_slide(blank_layout)
    add_header(slide9, "Hologram Keyboard & Screenshot (Label 3 & 4)")
    add_footer(slide9, 9)
    
    box1 = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, PInches(0.8), PInches(1.6), PInches(11.7), PInches(2.6))
    box1.fill.solid()
    box1.fill.fore_color.rgb = PRGBColor(240, 246, 252)
    box1.line.color.rgb = PRGBColor(0, 102, 204)
    tf1 = box1.text_frame
    tf1.word_wrap = True
    p = tf1.paragraphs[0]
    p.text = "LABEL 3 : PEACE / DUAL V-SIGN (Hologram Keyboard Summoning)"
    p.font.bold = True
    p.font.size = PPt(18)
    p.font.color.rgb = PRGBColor(0, 51, 102)
    p_sub = tf1.add_paragraph()
    p_sub.text = "• Summoning: Tahan gestur Peace (V-Sign) dengan kedua tangan selama 2 detik untuk memunculkan Keyboard Virtual Transparan (3 baris QWERTY).\n• Hover Typing: Mengetik dilakukan dengan mengarahkan (hover) jari telunjuk ke kotak tombol selama 1.0 detik. Bar progres hijau terisi penuh sebelum karakter diketik. Tombol khusus: Space, Backspace, Enter, & Close.\n• Keamanan: Saat keyboard aktif, navigasi kursor mouse dibekukan otomatis mencegah salah klik."
    p_sub.font.size = PPt(14)
    p_sub.font.color.rgb = PRGBColor(50, 50, 50)
    p_sub.space_before = PPt(6)

    box2 = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, PInches(0.8), PInches(4.4), PInches(11.7), PInches(2.5))
    box2.fill.solid()
    box2.fill.fore_color.rgb = PRGBColor(240, 246, 252)
    box2.line.color.rgb = PRGBColor(0, 102, 204)
    tf2 = box2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "LABEL 4 : FRAME KAMERA (Screenshot Area Mode)"
    p.font.bold = True
    p.font.size = PPt(18)
    p.font.color.rgb = PRGBColor(0, 51, 102)
    p_sub = tf2.add_paragraph()
    p_sub.text = "• Bounding Box Calculation: Membentuk bingkai kotak menggunakan jempol dan telunjuk dari kedua tangan. Sistem menghitung koordinat Bounding Box (x_min, y_min, x_max, y_max).\n• Eksekusi: Jika bingkai ditahan selama 1.5 detik (terdapat visual countdown bar & bingkai sudut hijau), sistem memetakan koordinat kamera ke resolusi monitor, memotong screenshot, menyimpannya ke Data/screenshots/, dan memberi efek kamera flash putih."
    p_sub.font.size = PPt(14)
    p_sub.font.color.rgb = PRGBColor(50, 50, 50)
    p_sub.space_before = PPt(6)

    # --------------------------------------------------------------------------
    # SLIDE 10: MULTITASKING & STANDBY LOCK (LABEL 6 & 7)
    # --------------------------------------------------------------------------
    slide10 = prs.slides.add_slide(blank_layout)
    add_header(slide10, "Task Switcher & Standby Lock (Label 6 & 7)")
    add_footer(slide10, 10)
    
    box1 = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, PInches(0.8), PInches(1.6), PInches(11.7), PInches(2.5))
    box1.fill.solid()
    box1.fill.fore_color.rgb = PRGBColor(240, 246, 252)
    box1.line.color.rgb = PRGBColor(0, 102, 204)
    tf1 = box1.text_frame
    tf1.word_wrap = True
    p = tf1.paragraphs[0]
    p.text = "LABEL 6 : TANGAN TERBUKA KERTAS (Swipe Tab / Task Switcher)"
    p.font.bold = True
    p.font.size = PPt(18)
    p.font.color.rgb = PRGBColor(0, 51, 102)
    p_sub = tf1.add_paragraph()
    p_sub.text = "• Kalkulasi Velocity (dx/dt): Eksklusif Tangan Kiri. Sistem mencatat history posisi X selama beberapa frame untuk menghitung kecepatan geser.\n• Aksi Alt+Tab: Swipe Cepat ke Kanan (Velocity > 1.3) memicu shortcut Alt + Tab (Next Tab). Swipe Cepat ke Kiri (Velocity < -1.3) memicu Alt + Shift + Tab (Prev Tab).\n• Cooldown: Dilengkapi cooldown 1.2 detik agar tidak terjadi perpindahan tab tak terkendali."
    p_sub.font.size = PPt(15)
    p_sub.font.color.rgb = PRGBColor(50, 50, 50)
    p_sub.space_before = PPt(8)

    box2 = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, PInches(0.8), PInches(4.4), PInches(11.7), PInches(2.5))
    box2.fill.solid()
    box2.fill.fore_color.rgb = PRGBColor(240, 246, 252)
    box2.line.color.rgb = PRGBColor(0, 102, 204)
    tf2 = box2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "LABEL 7 : TANGAN MENGEPAL KUAT / FIST (Standby / Pause Lock Mode)"
    p.font.bold = True
    p.font.size = PPt(18)
    p.font.color.rgb = PRGBColor(0, 51, 102)
    p_sub = tf2.add_paragraph()
    p_sub.text = "• System Lock: Genggam tangan (fist) selama 1.0 detik untuk mengunci sistem saat beristirahat. Seluruh kontrol mouse dan keyboard dimatikan, digantikan oleh overlay merah semi-transparan.\n• Unlock: Untuk membuka kunci, pengguna cukup menggenggam tangan kembali selama 1.0 detik.\n• Keamanan: Fitur ini mencegah kursor bergerak liar saat pengguna minum atau duduk santai."
    p_sub.font.size = PPt(15)
    p_sub.font.color.rgb = PRGBColor(50, 50, 50)
    p_sub.space_before = PPt(8)

    # --------------------------------------------------------------------------
    # SLIDE 11: 5 KEUNGGULAN TEKNIS DACTYL-AI
    # --------------------------------------------------------------------------
    slide11 = prs.slides.add_slide(blank_layout)
    add_header(slide11, "5 Keunggulan Teknis & Arsitektural Dactyl-AI")
    add_footer(slide11, 11)
    
    advs = [
        ("1. Dual-Wielding Bimanual Task Division", "Pembagian tugas kanan-kiri menghilangkan oklusi kamera dan konflik gestur."),
        ("2. Heuristic Jitter Filtering (EMA)", "Filter Exponential Moving Average (alpha=0.25) meredam tremor alami tangan secara real-time."),
        ("3. Lightweight AI Engine (1D Vector ML)", "SVM pada vektor numerik 126D jauh lebih ringan dibanding pemrosesan citra 2D CNN."),
        ("4. Non-Intrusive Always-on-Top GUI", "Win32 API membuat window kamera selalu di atas (topmost), transparan 60%, dan di pojok layar."),
        ("5. Asynchronous PowerShell Threading", "Eksekusi WMI Windows di background thread menjaga FPS kamera stabil di 24-30 FPS.")
    ]
    for idx, (title, desc) in enumerate(advs):
        top = PInches(1.5 + idx * 1.05)
        box = slide11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PInches(0.8), top, PInches(11.7), PInches(0.9))
        box.fill.solid()
        box.fill.fore_color.rgb = PRGBColor(240, 248, 255)
        box.line.color.rgb = PRGBColor(0, 102, 204)
        
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r1 = p.add_run()
        r1.text = f"{title} — "
        r1.font.bold = True
        r1.font.size = PPt(16)
        r1.font.color.rgb = PRGBColor(0, 51, 102)
        
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = PPt(14)
        r2.font.color.rgb = PRGBColor(60, 60, 60)

    # --------------------------------------------------------------------------
    # SLIDE 12: STRUKTUR KODE PROGRAM & PELATIHAN AI
    # --------------------------------------------------------------------------
    slide12 = prs.slides.add_slide(blank_layout)
    add_header(slide12, "Struktur Kode Program & Strategi Pelatihan AI")
    add_footer(slide12, 12)
    
    files = [
        ("data_collection.py", "Perekaman Dataset Bimanual", "Mengambil 126 koordinat dari MediaPipe dan menyimpannya ke Data/dataset_bimanual.csv. Perekaman dikontrol dengan menahan tombol keyboard [0] sampai [7] saat berpose."),
        ("train_model.py", "Pelatihan AI Brain (SVM)", "Membagi dataset menjadi 80% Train, 10% Validation, dan 10% Test dengan stratified sampling. Melakukan normalisasi StandardScaler dan melatih SVM Kernel RBF (C=10). Menyimpan model ke models/svm_gesture_model.pkl."),
        ("run_system.py", "Eksekusi Real-Time Pipeline", "Modul utama yang memuat model AI & scaler, menjalankan loop kamera 24-30 FPS, mendeteksi landmark, memprediksi gestur, dan mengeksekusi logika kontrol kursor, keyboard, & UI hologram.")
    ]
    for i, (fname, ftitle, fdesc) in enumerate(files):
        top = PInches(1.6 + i * 1.75)
        box = slide12.shapes.add_shape(MSO_SHAPE.RECTANGLE, PInches(0.8), top, PInches(11.7), PInches(1.55))
        box.fill.solid()
        box.fill.fore_color.rgb = PRGBColor(245, 248, 252)
        box.line.color.rgb = PRGBColor(0, 51, 102)
        box.line.width = PPt(1.5)
        
        tf = box.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = f"src/{fname}  |  {ftitle.upper()}"
        p1.font.bold = True
        p1.font.size = PPt(16)
        p1.font.color.rgb = PRGBColor(0, 51, 102)
        
        p2 = tf.add_paragraph()
        p2.text = fdesc
        p2.font.size = PPt(14)
        p2.font.color.rgb = PRGBColor(70, 70, 70)
        p2.space_before = PPt(6)

    # --------------------------------------------------------------------------
    # SLIDE 13: KESIMPULAN
    # --------------------------------------------------------------------------
    slide13 = prs.slides.add_slide(blank_layout)
    add_header(slide13, "Kesimpulan & Harapan Pengembangan")
    add_footer(slide13, 13)
    
    txBox = slide13.shapes.add_textbox(PInches(1.0), PInches(1.8), PInches(11.3), PInches(4.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    concl = [
        ("Sinergi AI & Geometri Spasial", "Dactyl-AI membuktikan bahwa penggabungan pemrosesan numerik 126 atribut MediaPipe dengan Support Vector Machine dan logika geometri spasial mampu menghasilkan kontrol sistem tanpa sentuhan yang berpresisi tinggi."),
        ("Solusi Praktis & Stabil", "Dengan performa real-time (24-30 FPS), fitur peredam getaran (EMA Jitter Filter), dan desain antarmuka Win32 transparan yang ergonomis, sistem ini menjadi solusi HCI yang sangat praktis dan stabil."),
        ("Masa Depan Antarmuka Komputer", "Dactyl-AI membuka jalan bagi generasi baru antarmuka pengguna di bidang realitas bertambah (AR/VR), ruang operasi medis steril, presentasi interaktif, hingga aksesibilitas bagi pengguna dengan keterbatasan fisik.")
    ]
    for idx, (head, body) in enumerate(concl):
        p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
        r1 = p.add_run()
        r1.text = f"✔ {head}\n"
        r1.font.bold = True
        r1.font.size = PPt(20)
        r1.font.color.rgb = PRGBColor(0, 51, 102)
        
        r2 = p.add_run()
        r2.text = body
        r2.font.size = PPt(16)
        r2.font.color.rgb = PRGBColor(60, 60, 60)
        p.space_before = PPt(10) if idx > 0 else PPt(0)
        p.space_after = PPt(15)

    prs.save(filename)
    print(f"[BERHASIL] Presentasi PowerPoint disimpan di: {os.path.abspath(filename)}")

if __name__ == "__main__":
    print("======================================================================")
    print("  MEMULAI PEMBUATAN DOKUMEN DAN PRESENTASI DACTYL-AI V2.5")
    print("======================================================================")
    
    output_dir = r"C:\Project Folder\AI\Dactyl-Ai\Dactyl-Ai"
    os.makedirs(output_dir, exist_ok=True)
    
    docx_path = os.path.join(output_dir, "Dactyl_AI_Dokumentasi_Lengkap.docx")
    pptx_path = os.path.join(output_dir, "Dactyl_AI_Presentasi_Lengkap.pptx")
    
    create_word_document(docx_path)
    create_powerpoint_presentation(pptx_path)
    print("======================================================================")
    print("  SELESAI! Semua file Word dan PowerPoint siap digunakan.")
    print("======================================================================")
